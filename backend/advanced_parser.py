"""
advanced_parser.py — Multi-format document ingestion with table + image extraction.

Extends the original parser.py (txt/vtt/pdf-text-only) to support the file
types a "general document" workflow actually needs — hiring brochures,
policy PDFs, Word docs, slide decks, spreadsheets — and to pull out
structured tables and image content (via OCR) rather than just flattening
everything into plain text and dropping the rest on the floor.

Supported extensions
---------------------
  .txt, .vtt            → delegated to parser.py (unchanged, meeting-style)
  .pdf                   → text (pypdf) + tables (pdfplumber) + images (OCR)
  .docx                  → paragraphs + tables + embedded images (OCR)
  .pptx                  → slide text + slide tables + slide images (OCR)
  .xlsx / .xls           → every sheet as a table (openpyxl)

Public API
----------
  is_supported(filename) -> bool
  extract(filename, raw_bytes) -> DocumentBundle

DocumentBundle is a dict:
  {
    "raw_text": str,                  # flattened readable text (for chat/extract context)
    "segments": list[dict],           # same shape parser.py uses: {speaker, text, timestamp}
    "tables":   list[dict],           # [{ "source": str, "caption": str|None, "headers": [...], "rows": [[...]] }]
    "images":   list[dict],           # [{ "source": str, "index": int, "ocr_text": str, "width": int, "height": int }]
  }

Design notes
------------
- OCR is best-effort: if tesseract isn't installed, images are still counted
  and reported, just without extracted text — the rest of the pipeline
  degrades gracefully rather than failing the whole upload.
- Table extraction never throws; a malformed/unextractable table is skipped
  and logged rather than blocking the whole document.
- Everything (tables + OCR'd image text) is folded into `raw_text` too, with
  clear markers ("[TABLE: ...]", "[IMAGE TEXT: ...]") so the existing
  chatbot / extractor prompt-builders — which only look at raw_text /
  segments — automatically pick this content up with zero changes to them.
"""

import io
import re
import logging
from typing import Optional

logger = logging.getLogger(__name__)

# ── Optional dependency imports (never hard-fail the module) ────────────────

try:
    import pdfplumber
    _PDFPLUMBER_OK = True
except ImportError:
    _PDFPLUMBER_OK = False

try:
    from pypdf import PdfReader
    _PYPDF_OK = True
except ImportError:
    _PYPDF_OK = False

try:
    import docx as _docx  # python-docx
    _DOCX_OK = True
except ImportError:
    _DOCX_OK = False

try:
    from pptx import Presentation
    _PPTX_OK = True
except ImportError:
    _PPTX_OK = False

try:
    import openpyxl
    _XLSX_OK = True
except ImportError:
    _XLSX_OK = False

try:
    from PIL import Image
    _PIL_OK = True
except ImportError:
    _PIL_OK = False

try:
    import pytesseract
    _OCR_OK = True
except ImportError:
    _OCR_OK = False


SUPPORTED_EXTENSIONS = ("txt", "vtt", "pdf", "docx", "pptx", "xlsx", "xls")
MAX_IMAGES_PER_DOC = 40     # OCR budget guard for very image-heavy files
MAX_OCR_CHARS_PER_IMAGE = 2000


def is_supported(filename: str) -> bool:
    ext = filename.rsplit(".", 1)[-1].lower() if "." in filename else ""
    return ext in SUPPORTED_EXTENSIONS


def _ocr_image_bytes(raw: bytes) -> tuple[str, int, int]:
    """Best-effort OCR. Returns (text, width, height); text is '' if OCR unavailable/failed."""
    if not _PIL_OK:
        return "", 0, 0
    try:
        img = Image.open(io.BytesIO(raw))
        w, h = img.size
    except Exception:
        return "", 0, 0

    if not _OCR_OK:
        return "", w, h

    try:
        # Skip tiny images (icons, bullets, decorative dividers) — not worth OCR
        # and a common source of OCR noise.
        if w < 60 or h < 60:
            return "", w, h
        text = pytesseract.image_to_string(img)
        text = text.strip()
        if len(text) > MAX_OCR_CHARS_PER_IMAGE:
            text = text[:MAX_OCR_CHARS_PER_IMAGE] + " …[truncated]"
        return text, w, h
    except Exception as exc:
        logger.warning(f"[AdvancedParser] OCR failed on an image: {exc}")
        return "", w, h


# ── PDF ───────────────────────────────────────────────────────────────────

def _extract_pdf(raw_bytes: bytes) -> dict:
    if not _PYPDF_OK:
        raise ValueError("PDF support requires the 'pypdf' package.")

    text_parts: list[str] = []
    tables: list[dict] = []
    images: list[dict] = []

    # --- Text (pypdf, same as before) ---
    try:
        reader = PdfReader(io.BytesIO(raw_bytes))
        if reader.is_encrypted:
            try:
                reader.decrypt("")
            except Exception:
                raise ValueError("This PDF is password-protected and cannot be read.")
        for page_num, page in enumerate(reader.pages, 1):
            try:
                t = page.extract_text() or ""
            except Exception:
                t = ""
            if t.strip():
                text_parts.append(t.strip())

            # --- Embedded images (pypdf) ---
            if len(images) < MAX_IMAGES_PER_DOC:
                try:
                    for img_idx, img_file in enumerate(page.images, 1):
                        if len(images) >= MAX_IMAGES_PER_DOC:
                            break
                        ocr_text, w, h = _ocr_image_bytes(img_file.data)
                        images.append({
                            "source": f"page {page_num}",
                            "index": img_idx,
                            "ocr_text": ocr_text,
                            "width": w, "height": h,
                        })
                except Exception as exc:
                    logger.warning(f"[AdvancedParser] Image extraction failed on page {page_num}: {exc}")
    except ValueError:
        raise
    except Exception as exc:
        raise ValueError(f"Could not read PDF file: {exc}") from exc

    # --- Tables (pdfplumber) ---
    if _PDFPLUMBER_OK:
        try:
            with pdfplumber.open(io.BytesIO(raw_bytes)) as pdf:
                for page_num, page in enumerate(pdf.pages, 1):
                    try:
                        found = page.extract_tables()
                    except Exception:
                        found = []
                    for t_idx, raw_table in enumerate(found, 1):
                        table = _clean_table(raw_table)
                        if table:
                            tables.append({
                                "source": f"page {page_num}, table {t_idx}",
                                "caption": None,
                                "headers": table["headers"],
                                "rows": table["rows"],
                            })
        except Exception as exc:
            logger.warning(f"[AdvancedParser] Table extraction failed: {exc}")

    full_text = "\n\n".join(text_parts).strip()
    if not full_text and not tables and not images:
        raise ValueError(
            "No extractable content found in this PDF. It may be a scanned/image-only "
            "document without a text layer — try exporting a .txt or .vtt transcript instead."
        )
    return {"raw_text": full_text, "tables": tables, "images": images}


# ── DOCX ──────────────────────────────────────────────────────────────────

def _extract_docx(raw_bytes: bytes) -> dict:
    if not _DOCX_OK:
        raise ValueError("Word document support requires the 'python-docx' package.")

    try:
        doc = _docx.Document(io.BytesIO(raw_bytes))
    except Exception as exc:
        raise ValueError(f"Could not read Word document: {exc}") from exc

    text_parts = []
    for para in doc.paragraphs:
        if para.text.strip():
            # Preserve heading structure lightly so section context survives.
            style = (para.style.name or "").lower()
            prefix = "## " if "heading" in style else ""
            text_parts.append(prefix + para.text.strip())

    tables = []
    for t_idx, table in enumerate(doc.tables, 1):
        rows = [[cell.text.strip() for cell in row.cells] for row in table.rows]
        cleaned = _clean_table(rows)
        if cleaned:
            tables.append({
                "source": f"table {t_idx}",
                "caption": None,
                "headers": cleaned["headers"],
                "rows": cleaned["rows"],
            })

    images = []
    try:
        for img_idx, rel in enumerate(
            [r for r in doc.part.rels.values() if "image" in r.reltype], 1
        ):
            if img_idx > MAX_IMAGES_PER_DOC:
                break
            try:
                blob = rel.target_part.blob
            except Exception:
                continue
            ocr_text, w, h = _ocr_image_bytes(blob)
            images.append({
                "source": "document", "index": img_idx,
                "ocr_text": ocr_text, "width": w, "height": h,
            })
    except Exception as exc:
        logger.warning(f"[AdvancedParser] DOCX image extraction failed: {exc}")

    full_text = "\n".join(text_parts).strip()
    if not full_text and not tables and not images:
        raise ValueError("No extractable content found in this Word document.")
    return {"raw_text": full_text, "tables": tables, "images": images}


# ── PPTX ──────────────────────────────────────────────────────────────────

def _extract_pptx(raw_bytes: bytes) -> dict:
    if not _PPTX_OK:
        raise ValueError("PowerPoint support requires the 'python-pptx' package.")

    try:
        prs = Presentation(io.BytesIO(raw_bytes))
    except Exception as exc:
        raise ValueError(f"Could not read PowerPoint file: {exc}") from exc

    text_parts = []
    tables = []
    images = []

    for slide_num, slide in enumerate(prs.slides, 1):
        slide_lines = [f"## Slide {slide_num}"]
        for shape in slide.shapes:
            if shape.has_text_frame and shape.text_frame.text.strip():
                slide_lines.append(shape.text_frame.text.strip())

            if shape.has_table:
                raw_table = [[cell.text.strip() for cell in row.cells] for row in shape.table.rows]
                cleaned = _clean_table(raw_table)
                if cleaned:
                    tables.append({
                        "source": f"slide {slide_num}",
                        "caption": None,
                        "headers": cleaned["headers"],
                        "rows": cleaned["rows"],
                    })

            if shape.shape_type == 13 and len(images) < MAX_IMAGES_PER_DOC:  # 13 = PICTURE
                try:
                    blob = shape.image.blob
                    ocr_text, w, h = _ocr_image_bytes(blob)
                    images.append({
                        "source": f"slide {slide_num}", "index": len(images) + 1,
                        "ocr_text": ocr_text, "width": w, "height": h,
                    })
                except Exception:
                    pass

        if len(slide_lines) > 1:
            text_parts.append("\n".join(slide_lines))

    full_text = "\n\n".join(text_parts).strip()
    if not full_text and not tables and not images:
        raise ValueError("No extractable content found in this PowerPoint file.")
    return {"raw_text": full_text, "tables": tables, "images": images}


# ── XLSX ──────────────────────────────────────────────────────────────────

def _extract_xlsx(raw_bytes: bytes) -> dict:
    if not _XLSX_OK:
        raise ValueError("Excel support requires the 'openpyxl' package.")

    try:
        wb = openpyxl.load_workbook(io.BytesIO(raw_bytes), data_only=True, read_only=True)
    except Exception as exc:
        raise ValueError(f"Could not read Excel file: {exc}") from exc

    tables = []
    text_parts = []
    for sheet in wb.worksheets:
        rows = []
        for row in sheet.iter_rows(values_only=True):
            if any(c is not None for c in row):
                rows.append(["" if c is None else str(c) for c in row])
            if len(rows) > 2000:  # sanity cap per sheet
                break
        cleaned = _clean_table(rows)
        if cleaned:
            tables.append({
                "source": f"sheet '{sheet.title}'",
                "caption": sheet.title,
                "headers": cleaned["headers"],
                "rows": cleaned["rows"],
            })
            text_parts.append(f"## Sheet: {sheet.title}")

    full_text = "\n".join(text_parts).strip()
    if not tables:
        raise ValueError("No data found in this spreadsheet.")
    return {"raw_text": full_text, "tables": tables, "images": []}


# ── Table cleanup ─────────────────────────────────────────────────────────

def _clean_table(rows: list[list]) -> Optional[dict]:
    """Drop empty rows/cols, assume first row is a header if it looks like one."""
    if not rows:
        return None
    rows = [[("" if c is None else str(c).strip()) for c in r] for r in rows]
    rows = [r for r in rows if any(c for c in r)]
    if len(rows) < 1:
        return None
    # Trim to the widest non-empty row for consistency
    width = max(len(r) for r in rows)
    rows = [r + [""] * (width - len(r)) for r in rows]

    headers = rows[0] if len(rows) > 1 else [f"col_{i+1}" for i in range(width)]
    body = rows[1:] if len(rows) > 1 else rows
    if not body:
        return None
    return {"headers": headers, "rows": body}


def table_to_markdown(table: dict, max_rows: int = 30) -> str:
    """Render a table dict as a markdown table string (for LLM context)."""
    headers = table.get("headers") or []
    rows = table.get("rows") or []
    if not headers or not rows:
        return ""
    lines = [
        "| " + " | ".join(headers) + " |",
        "| " + " | ".join(["---"] * len(headers)) + " |",
    ]
    for r in rows[:max_rows]:
        lines.append("| " + " | ".join(str(c) for c in r) + " |")
    if len(rows) > max_rows:
        lines.append(f"_(+{len(rows) - max_rows} more rows omitted)_")
    return "\n".join(lines)


# ── Public entry point ──────────────────────────────────────────────────────

def extract(filename: str, raw_bytes: bytes) -> dict:
    """
    Dispatch by extension. Returns a DocumentBundle:
      { "raw_text": str, "tables": [...], "images": [...] }
    Raises ValueError with a user-facing message on failure.
    """
    ext = filename.rsplit(".", 1)[-1].lower() if "." in filename else ""

    if ext == "pdf":
        bundle = _extract_pdf(raw_bytes)
    elif ext == "docx":
        bundle = _extract_docx(raw_bytes)
    elif ext == "pptx":
        bundle = _extract_pptx(raw_bytes)
    elif ext in ("xlsx", "xls"):
        bundle = _extract_xlsx(raw_bytes)
    else:
        raise ValueError(f"Unsupported file type '.{ext}' for advanced parsing.")

    # Fold tables + OCR'd image text into raw_text so every existing
    # consumer (chatbot context builder, extractor) sees this content
    # automatically, with zero changes needed on their end.
    enriched = [bundle["raw_text"]] if bundle["raw_text"] else []

    for table in bundle["tables"]:
        md = table_to_markdown(table)
        if md:
            enriched.append(f"[TABLE — {table['source']}]\n{md}")

    for img in bundle["images"]:
        if img.get("ocr_text"):
            enriched.append(f"[IMAGE TEXT — {img['source']}]\n{img['ocr_text']}")

    bundle["raw_text"] = "\n\n".join(enriched).strip()
    return bundle
